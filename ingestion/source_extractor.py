"""源文档导入与结构化提取模块

从 MySQL 导出的 SQL dump 导入 source_documents 表，
按 file_type 分发提取逻辑，提取后灌入 raw_items 进入 cleaning pipeline。

Extract+Clean 合并架构：
  _extract_single() 纯提取 → _semantic_clean() DeepSeek 清洗/理解
  根据 file_type 和是否过了 Qwen 视觉模型，选择不同的清洗策略：
  - 过了 Qwen 的 (PDF图表描述、image/mixed 视觉降级) → 轻度清洗（去格式噪音）
  - 没过 Qwen 的 OCR/裸文本 → 深度理解整理（重新组织信息结构）
"""
import json
import logging
import os
import re
import tempfile
import threading

from utils.db_utils import execute_cloud_query, execute_cloud_insert, get_db

logger = logging.getLogger(__name__)

# ==================== DeepSeek 清洗客户端（lazy singleton）====================

_deepseek_client = None
_deepseek_lock = threading.Lock()


def _get_deepseek():
    global _deepseek_client
    if _deepseek_client is None:
        with _deepseek_lock:
            if _deepseek_client is None:
                from openai import OpenAI
                rows = execute_cloud_query(
                    "SELECT value FROM system_config WHERE config_key='deepseek_api_key'"
                )
                if not rows:
                    raise RuntimeError("system_config 中未找到 deepseek_api_key")
                _deepseek_client = OpenAI(
                    api_key=rows[0]["value"], base_url="https://api.deepseek.com/v1"
                )
    return _deepseek_client


def _call_deepseek(system_prompt: str, text: str, max_tokens=4096, timeout=90) -> str:
    if len(text) > 12000:
        text = text[:12000] + "\n\n[文本已截断]"
    client = _get_deepseek()
    resp = client.chat.completions.create(
        model="deepseek-chat",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": text},
        ],
        max_tokens=max_tokens,
        temperature=0.3,
        timeout=timeout,
    )
    return resp.choices[0].message.content


# ==================== 清洗 Prompts ====================

# --- 理解整理型（没过 Qwen 的 OCR / 裸文本） ---
_UNDERSTAND_PROMPT = """你是金融信息整理助手。以下文本从截图OCR/帖子中提取，包含大量UI残留、导航文字、按钮、乱码等噪音。

请执行：
1. 去掉所有UI元素（搜索框、导航栏、按钮、页码、网速显示、App界面文字等）
2. 去掉乱码和无意义字符
3. 理解内容后重新组织成结构清晰的信息摘要
4. 保留所有关键数据（涨跌幅、市占率、产能数字、公司名、股票代码等）
5. 如果是表格类数据，用 Markdown 表格格式整理
6. 如果文本已经很干净且结构清晰，直接原样返回

直接输出整理后的文本，不要加任何前缀说明。"""

# --- 轻度清洗型（已过 Qwen 或文本本身较干净） ---
_CLEAN_PROMPTS = {
    "pdf": """你是金融研报文本清洗专家。以下文本从PDF提取，由于多栏布局，侧边栏信息（分析师信息、股价评级、表现数据、免责声明等）可能被混入正文中间，打断了原本连贯的句子。

请执行：
1. 找到所有打断正文语义的异物文字（通常出现在句子中间，与前后文不连贯）
2. 删除这些异物文字
3. 修复被打断的句子，使其恢复连贯
4. 删除页眉页脚、页码、免责声明等重复出现的模板文字
5. 保留所有正文内容、数据、表格、图表描述、目录

如果文本已经很干净，直接原样返回。只做删除和修复，不添加任何新内容。直接输出清洗后的文本。""",

    "audio": """你是金融电话会议/音频转写文本清洗专家。以下文本由语音识别转写而来，包含口语噪音。

请执行：
1. 删除无意义的口语填充词和重复（"嗯"、"啊"、"那个"、"就是说"、连续重复的词句）
2. 修复被打断的句子——电话会中常有人插话打断，导致一句话被切成两段夹着别人的话
3. 合并说话人的断续表达，使其成为完整连贯的句子
4. 删除主持人的程序性话术（"下面有请XX回答"、"感谢XX的提问"等）
5. 保留所有实质性内容：观点、数据、问答、业务讨论
6. 保留说话人标识（如有）

不要改变原意，不要添加新内容，不要总结。直接输出清洗后的文本。""",
}
_CLEAN_PROMPTS["mp3"] = _CLEAN_PROMPTS["audio"]


def _split_long_text(text: str, chunk_size: int = 5000, overlap: int = 300) -> list:
    """长文本分段，优先按自然段落切分"""
    if len(text) <= chunk_size:
        return [text]
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size, chunk_overlap=overlap,
            separators=["\n\n", "\n", "。", "；", " "],
        )
        return splitter.split_text(text)
    except ImportError:
        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunks.append(text[i:i + chunk_size])
        return chunks


def _semantic_clean(raw_text: str, file_type: str, doc_id, needs_understanding: bool = False) -> str:
    """Extract+Clean 合并清洗入口

    Args:
        raw_text: 提取后的原始文本
        file_type: 文件类型
        doc_id: source_documents.id（日志用）
        needs_understanding: True=没过Qwen，需要深度理解整理；False=轻度清洗
    Returns:
        清洗后的文本
    """
    if not raw_text or len(raw_text.strip()) < 50:
        return raw_text

    # xlsx 不清洗
    if file_type in ("xlsx", "xls"):
        return raw_text

    # 选择 prompt
    if needs_understanding:
        prompt = _UNDERSTAND_PROMPT
    elif file_type in _CLEAN_PROMPTS:
        prompt = _CLEAN_PROMPTS[file_type]
    else:
        # 兜底：理解整理
        prompt = _UNDERSTAND_PROMPT

    chunks = _split_long_text(raw_text)
    cleaned_parts = []
    for i, chunk in enumerate(chunks):
        try:
            result = _call_deepseek(prompt, chunk, max_tokens=4096, timeout=60)
            cleaned_parts.append(result if result else chunk)
        except Exception as e:
            logger.warning(f"[S] 语义清洗第{i+1}/{len(chunks)}段失败 doc_id={doc_id}: {e}")
            cleaned_parts.append(chunk)

    cleaned = "\n\n".join(cleaned_parts)

    if cleaned and len(cleaned) > 50:
        logger.info(
            f"[S] 语义清洗完成 doc_id={doc_id} type={file_type} "
            f"understand={needs_understanding} ({len(raw_text)}→{len(cleaned)}字, {len(chunks)}段)"
        )
        return cleaned

    logger.info(f"[S] 语义清洗结果过短，保留原文 doc_id={doc_id}")
    return raw_text


# ==================== SQL 导入 ====================

def import_sql_dump(sql_path: str) -> dict:
    """解析 MySQL INSERT 语句并导入到 source_documents

    Args:
        sql_path: SQL 文件路径
    Returns:
        {"total": N, "imported": M, "skipped": K}
    """
    with open(sql_path, "r", encoding="utf-8") as f:
        sql_text = f.read()

    # 提取所有 INSERT 语句中的 VALUES 行
    # 格式: INSERT INTO `stock_analysis` VALUES (id, 'doc_type', 'file_type', ...)
    pattern = re.compile(
        r"INSERT\s+INTO\s+`stock_analysis`\s+VALUES\s*\((.+?)\);",
        re.DOTALL,
    )
    matches = pattern.findall(sql_text)

    total = len(matches)
    imported = 0
    skipped = 0

    for row_str in matches:
        try:
            fields = _parse_values_row(row_str)
            if len(fields) < 12:
                logger.warning(f"字段数不足，跳过: {row_str[:80]}...")
                skipped += 1
                continue

            doc_id = int(fields[0])
            doc_type = fields[1]
            file_type = fields[2]
            title = fields[3]
            author = fields[4]
            publish_date = fields[5]
            source = fields[6]
            oss_url = fields[7]
            text_content = fields[8]
            # fields[9] = status (skip)
            created_at = fields[10]
            updated_at = fields[11]

            execute_cloud_insert(
                """INSERT IGNORE INTO source_documents
                   (id, doc_type, file_type, title, author, publish_date,
                    source, oss_url, text_content, created_at, updated_at)
                   VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)""",
                [doc_id, doc_type, file_type, title, author, publish_date,
                 source if source != "NULL" else None,
                 oss_url if oss_url != "NULL" else None,
                 text_content if text_content != "NULL" else None,
                 created_at, updated_at],
            )
            imported += 1
        except Exception as e:
            logger.error(f"导入行失败: {e} — {row_str[:80]}...")
            skipped += 1

    logger.info(f"SQL导入完成: total={total}, imported={imported}, skipped={skipped}")
    return {"total": total, "imported": imported, "skipped": skipped}


def _parse_values_row(row_str: str) -> list:
    """解析单条 INSERT VALUES 中的字段列表

    处理嵌套引号和 NULL 值。
    """
    fields = []
    i = 0
    n = len(row_str)

    while i < n:
        # 跳过空白和逗号
        while i < n and row_str[i] in (" ", "\t", "\n", "\r", ","):
            i += 1
        if i >= n:
            break

        if row_str[i] == "'":
            # 带引号的字符串
            i += 1  # skip opening quote
            val_parts = []
            while i < n:
                if row_str[i] == "\\" and i + 1 < n:
                    # 转义字符
                    next_char = row_str[i + 1]
                    if next_char == "'":
                        val_parts.append("'")
                    elif next_char == "\\":
                        val_parts.append("\\")
                    elif next_char == "n":
                        val_parts.append("\n")
                    elif next_char == "r":
                        val_parts.append("\r")
                    elif next_char == "t":
                        val_parts.append("\t")
                    else:
                        val_parts.append(next_char)
                    i += 2
                elif row_str[i] == "'":
                    i += 1  # skip closing quote
                    break
                else:
                    val_parts.append(row_str[i])
                    i += 1
            fields.append("".join(val_parts))
        else:
            # 无引号值 (数字或 NULL)
            j = i
            while j < n and row_str[j] not in (",", ")"):
                j += 1
            val = row_str[i:j].strip()
            fields.append(val)
            i = j

    return fields


# ==================== 文本提取 ====================

def extract_batch(file_type: str = None, limit: int = 50, on_progress=None) -> dict:
    """批量提取源文档

    Args:
        file_type: 限定文件类型 (txt/image/mixed/pdf/audio), None=全部
        limit: 每次处理数
        on_progress(done, total, row_id): 每完成一条时回调
    Returns:
        {"total": N, "success": M, "failed": K}
    """
    sql = """SELECT id, doc_type, file_type, title, text_content, oss_url
             FROM source_documents WHERE extract_status='pending'"""
    params = []
    if file_type:
        sql += " AND file_type=%s"
        params.append(file_type)
    sql += " LIMIT %s"
    params.append(limit)

    rows = execute_cloud_query(sql, params)
    total = len(rows)
    success = 0
    failed = 0

    if on_progress:
        on_progress(0, total, None)

    for i, row in enumerate(rows):
        try:
            extracted = _extract_and_clean_single(row)
            from config.doc_types import classify_doc_type
            new_doc_type = classify_doc_type(
                row.get("title") or "",
                (extracted or "")[:200],
            )
            execute_cloud_insert(
                """UPDATE source_documents
                   SET extracted_text=%s, extract_status='extracted', doc_type=%s
                   WHERE id=%s""",
                [extracted, new_doc_type, row["id"]],
            )
            success += 1
        except Exception as e:
            err_str = str(e)
            # 401/403 表示 URL 已过期或无权限，跳过而不是标记失败
            if "401" in err_str or "403" in err_str or "Unauthorized" in err_str or "Forbidden" in err_str:
                logger.warning(f"跳过 id={row['id']} (URL 过期/无权限): {err_str[:100]}")
                execute_cloud_insert(
                    """UPDATE source_documents
                       SET extract_status='skipped'
                       WHERE id=%s""",
                    [row["id"]],
                )
            else:
                logger.error(f"提取失败 id={row['id']}: {e}")
                execute_cloud_insert(
                    """UPDATE source_documents
                       SET extract_status='failed'
                       WHERE id=%s""",
                    [row["id"]],
                )
            failed += 1
        if on_progress:
            on_progress(i + 1, total, row["id"])

    logger.info(f"批量提取完成: total={total}, success={success}, failed={failed}")
    return {"total": total, "success": success, "failed": failed}


def _extract_single(row: dict) -> str:
    """按 file_type 分发提取（纯提取，不含清洗）"""
    text, _ = _extract_single_with_meta(row)
    return text


def _extract_single_with_meta(row: dict) -> tuple:
    """按 file_type 分发提取，返回 (text, needs_understanding)

    needs_understanding=True 表示内容没过 Qwen 视觉理解，需要 DeepSeek 深度理解整理
    needs_understanding=False 表示已过 Qwen 或是结构化数据，轻度清洗即可
    """
    ft = row["file_type"]
    if ft == "txt":
        return _extract_txt(row), True  # 裸文本，需要理解整理
    elif ft == "mp3":
        return _extract_mp3(row), False  # 音频转写，用 audio 清洗 prompt
    elif ft == "image":
        return _extract_image_with_meta(row)  # 取决于是否走了 Qwen
    elif ft == "mixed":
        return _extract_mixed_with_meta(row)  # 取决于是否走了 Qwen
    elif ft == "pdf":
        return _extract_pdf_with_meta(row)  # 正常 PDF=False，扫描件=True
    elif ft == "audio":
        return _extract_audio(row), False  # 音频转写，用 audio 清洗 prompt
    elif ft in ("xlsx", "xls"):
        return _extract_xlsx(row), False  # 结构化数据，不清洗
    else:
        raise ValueError(f"未知 file_type: {ft}")


def _extract_and_clean_single(row: dict) -> str:
    """提取 + 清洗一体化：提取后立即做语义清洗"""
    text, needs_understanding = _extract_single_with_meta(row)
    if not text or len(text.strip()) < 20:
        return text
    return _semantic_clean(text, row["file_type"], row["id"], needs_understanding)


def _extract_txt(row: dict) -> str:
    """txt: 直接使用 text_content，并解析 ZSXQ 富文本标签"""
    import re
    from urllib.parse import unquote
    text = row.get("text_content") or ""

    def replace_tag(m):
        tag_type = m.group(1)
        attrs = m.group(2)
        if tag_type == "web":
            href = re.search(r'href="([^"]+)"', attrs)
            title = re.search(r'title="([^"]+)"', attrs)
            url = unquote(href.group(1)) if href else ""
            label = unquote(title.group(1)) if title else url
            return f"[{label}]({url})"
        elif tag_type == "hashtag":
            title = re.search(r'title="([^"]+)"', attrs)
            return unquote(title.group(1)) if title else ""
        elif tag_type == "mention":
            title = re.search(r'title="([^"]+)"', attrs)
            return f"@{unquote(title.group(1))}" if title else ""
        return ""

    text = re.sub(r'<e\s+type="(\w+)"([^/]*)/?>', replace_tag, text)
    return text


def _forward_fill(rows: list[list[str]]) -> list[list[str]]:
    """对二维表做列方向 forward-fill，填充合并单元格留下的空洞。

    只填充非标题行（跳过第一行），且只填充那些在首行有值的列，
    避免把真正的空值也填上。
    """
    if len(rows) < 2:
        return rows

    # 判断哪些列需要 forward-fill：首行有值的列才可能是行标签
    header = rows[0]
    fill_cols = {i for i, v in enumerate(header) if v.strip()}

    last = list(rows[0])
    result = [rows[0]]
    for row in rows[1:]:
        filled = list(row)
        for i in fill_cols:
            if i < len(filled) and filled[i] == "" and i < len(last):
                filled[i] = last[i]
            if i < len(filled) and filled[i] != "":
                last[i] = filled[i]
        result.append(filled)
    return result


def _extract_xlsx(row: dict) -> str:
    """xlsx/xls: 全量提取所有 sheet 所有行转 CSV，并用 LLM 对前50行生成摘要用于索引"""
    import io
    import csv
    import openpyxl

    oss_url = row.get("oss_url") or ""
    if not oss_url:
        return row.get("text_content") or ""

    data = _download_file(oss_url, timeout=60)

    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    except Exception as e:
        logger.warning(f"openpyxl 打开失败: {e}")
        return row.get("text_content") or ""

    sheet_csvs = []      # 全量 CSV（每个 sheet）
    sheet_previews = []  # 前50行预览（用于 LLM 摘要）

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_rows = []
        for row_cells in ws.iter_rows(values_only=True):
            if not any(v is not None and str(v).strip() for v in row_cells):
                continue
            all_rows.append([str(v) if v is not None else "" for v in row_cells])

        if not all_rows:
            continue

        all_rows = _forward_fill(all_rows)

        buf = io.StringIO()
        csv.writer(buf).writerows(all_rows)
        sheet_csvs.append(f"## Sheet: {sheet_name}\n{buf.getvalue().strip()}")

        # 前50行用于摘要
        preview_buf = io.StringIO()
        csv.writer(preview_buf).writerows(all_rows[:50])
        sheet_previews.append(f"## Sheet: {sheet_name}\n{preview_buf.getvalue().strip()}")

    wb.close()

    if not sheet_csvs:
        return row.get("text_content") or ""

    full_csv = "\n\n".join(sheet_csvs)
    preview_text = "\n\n".join(sheet_previews)

    summary = _summarize_xlsx_with_llm(row.get("title", ""), preview_text)

    if summary:
        return f"【摘要】\n{summary}\n\n---\n【完整数据】\n{full_csv}"
    return full_csv


def _summarize_xlsx_with_llm(title: str, preview_csv: str) -> str:
    """用 LLM 对表格前50行生成结构化摘要，用于后续索引和检索"""
    prompt = f"""以下是一份名为「{title}」的 Excel 表格数据（CSV格式，前50行）。

请分析并输出：
1. 数据概述（包含哪些指标/维度）
2. 关键统计特征（最大值、最小值、均值、变化幅度等）
3. 主要趋势或方向（上升/下降/分化等）
4. 重要占比或结构关系
5. 值得关注的异常或亮点

用简洁的中文输出，不超过400字。

数据：
{preview_csv[:6000]}"""

    try:
        from utils.model_router import call_model
        return call_model("cleaning", prompt) or ""
    except Exception as e:
        logger.warning(f"xlsx LLM 摘要失败: {e}")
        return ""


def _get_audio_config() -> dict:
    """读取 audio stage 配置

    返回 {'provider': str, 'model_name': str, 'groq_key': str, 'gemini_key': str}
    provider='funasr' 时走本地，无需 key。
    """
    from utils.db_utils import get_config
    result = {"provider": "", "model_name": "", "groq_key": "", "gemini_key": ""}

    try:
        from utils.model_router import _load_stage_config
        cfg = _load_stage_config("audio")
        provider = (cfg.get("provider") or "").lower()
        model_name = cfg.get("model_name") or ""
        api_key_ref = cfg.get("api_key_ref") or ""
        result["provider"] = provider
        result["model_name"] = model_name
        if api_key_ref:
            key_val = get_config(api_key_ref) or ""
            if provider == "groq":
                result["groq_key"] = key_val
            elif provider in ("gemini", "google"):
                result["gemini_key"] = key_val
            elif "groq" in api_key_ref:
                result["groq_key"] = key_val
            elif "gemini" in api_key_ref or "google" in api_key_ref:
                result["gemini_key"] = key_val
    except Exception:
        pass

    # 回退到环境变量
    if not result["groq_key"]:
        result["groq_key"] = os.getenv("GROQ_API_KEY", "")
    if not result["gemini_key"]:
        result["gemini_key"] = os.getenv("GEMINI_API_KEY", "")

    return result


def _transcribe_funasr(audio_path: str, model_name: str = "iic/SenseVoiceSmall") -> str:
    """用 FunASR 本地模型转写音频（无需 API key）"""
    from funasr import AutoModel
    model = AutoModel(
        model=model_name,
        vad_model="fsmn-vad",
        punc_model="ct-punc",
        disable_update=True,
    )
    res = model.generate(input=audio_path, batch_size_s=300)
    if not res:
        return ""
    text = res[0].get("text", "")
    # SenseVoice 输出带情绪标签如 <|zh|><|NEUTRAL|>，清理掉
    text = re.sub(r"<\|[^|]+\|>", "", text).strip()
    return text



def _extract_mp3(row: dict) -> str:
    """mp3: 下载 URL 后转写（兼容两种存储方式）

    旧格式: oss_url=下载URL, text_content=帖子原文
    新格式: text_content=下载URL
    优先级: FunASR（本地）→ Groq → Gemini
    """
    tc = (row.get("text_content") or "").strip()
    url = tc if tc.startswith("http") else (row.get("oss_url") or "").strip()
    if not url or not url.startswith("http"):
        return ""

    import time

    data = _download_file(url, timeout=120)

    with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        cfg = _get_audio_config()

        # FunASR 本地优先
        if cfg["provider"] == "funasr":
            try:
                return _transcribe_funasr(tmp_path, cfg["model_name"] or "iic/SenseVoiceSmall")
            except Exception as e:
                logger.warning(f"FunASR 失败，降级到 Groq: {e}")

        groq_key = cfg["groq_key"]
        gemini_key = cfg["gemini_key"]

        if groq_key:
            try:
                return _transcribe_groq(tmp_path, groq_key)
            except Exception as e:
                if "rate_limit" in str(e) or "429" in str(e):
                    logger.warning("Groq 速率限制，切换 Gemini...")
                else:
                    raise

        if gemini_key:
            try:
                return _transcribe_gemini(data, gemini_key)
            except Exception as e:
                if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                    logger.warning("Gemini 也受限...")
                else:
                    raise

        if groq_key:
            logger.warning("两个 API 都受限，等待 10 分钟后重试 Groq...")
            time.sleep(620)
            return _transcribe_groq(tmp_path, groq_key)

        raise RuntimeError("未配置任何音频转写方案（funasr/groq/gemini）")
    finally:
        os.unlink(tmp_path)




def _transcribe_groq(tmp_path: str, api_key: str) -> str:
    """用 Groq Whisper API 转写音频（带时间戳段落标记）"""
    from groq import Groq
    client = Groq(api_key=api_key)
    with open(tmp_path, "rb") as audio_file:
        try:
            transcription = client.audio.transcriptions.create(
                file=("audio.mp3", audio_file),
                model="whisper-large-v3",
                language="zh",
                response_format="verbose_json",
            )
            # verbose_json 有 segments[{start, end, text}]
            segments = getattr(transcription, "segments", None)
            if segments and len(segments) > 1:
                return _format_audio_segments(segments)
            return getattr(transcription, "text", str(transcription))
        except Exception:
            # 降级到普通模式
            audio_file.seek(0)
            transcription = client.audio.transcriptions.create(
                file=("audio.mp3", audio_file),
                model="whisper-large-v3",
                language="zh",
            )
            return transcription.text


def _format_audio_segments(segments, group_seconds: int = 300) -> str:
    """将音频转写段落按时间分组，每 group_seconds 秒一个段落标记"""
    parts = []
    current_group_start = 0
    current_texts = []

    for seg in segments:
        start = seg.get("start") or seg.start if hasattr(seg, "start") else 0
        text = seg.get("text") or seg.text if hasattr(seg, "text") else ""
        text = text.strip()
        if not text:
            continue

        if start >= current_group_start + group_seconds and current_texts:
            mins, secs = divmod(int(current_group_start), 60)
            hours, mins = divmod(mins, 60)
            time_label = f"{hours:02d}:{mins:02d}:{secs:02d}"
            parts.append(f"### [{time_label}]\n{''.join(current_texts)}")
            current_group_start = int(start // group_seconds) * group_seconds
            current_texts = []

        current_texts.append(text)

    # 最后一组
    if current_texts:
        mins, secs = divmod(int(current_group_start), 60)
        hours, mins = divmod(mins, 60)
        time_label = f"{hours:02d}:{mins:02d}:{secs:02d}"
        parts.append(f"### [{time_label}]\n{''.join(current_texts)}")

    return "\n\n".join(parts) if parts else "".join(t for seg in segments for t in [seg.get("text", "")])


def _transcribe_gemini(audio_data: bytes, api_key: str) -> str:
    """用 Gemini API 转写音频"""
    from google import genai
    from google.genai import types
    client = genai.Client(api_key=api_key)
    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=[
            types.Content(parts=[
                types.Part(text="请将这段音频完整转写为中文文字。只输出转写内容，不要添加任何额外说明。"),
                types.Part(inline_data=types.Blob(mime_type="audio/mpeg", data=audio_data)),
            ])
        ],
    )
    return response.text


def _extract_image(row: dict) -> str:
    """image: 解析 [图片N] url 格式（兼容旧调用）"""
    text, _ = _extract_image_with_meta(row)
    return text


def _extract_image_with_meta(row: dict) -> tuple:
    """image: 解析 [图片N] url 格式，返回 (text, needs_understanding)

    策略：PaddleOCR 提取文字 → 文字太少时降级到视觉模型描述
    needs_understanding=True 当所有图片都走了 OCR（没过 Qwen）
    needs_understanding=False 当有图片走了 Qwen 视觉模型
    """
    text_content = row.get("text_content") or ""
    urls = _parse_image_urls(text_content)
    if not urls:
        return text_content, True  # 无图片URL，纯文本，需要理解

    ocr_parts = []
    non_text_urls = []
    non_text_indices = []

    for idx, url in enumerate(urls):
        try:
            img_data = _download_file(url, timeout=30)
            ocr_text = _ocr_image_bytes(img_data)
            if len(ocr_text.strip()) >= 20:
                ocr_parts.append(f"### [图片{idx+1} OCR文字]\n{ocr_text}")
            else:
                non_text_urls.append(url)
                non_text_indices.append(idx)
        except Exception as e:
            logger.warning(f"PaddleOCR 图片失败: {e}")
            non_text_urls.append(url)
            non_text_indices.append(idx)

    result_parts = ocr_parts[:]
    used_qwen = False

    # 无文字的图片降级到视觉模型
    if non_text_urls:
        try:
            from utils.model_router import call_model_vision
            prompt = (
                f"这是一张关于'{row.get('title', '股票分析')}'的图片。"
                "请详细描述图片中的所有内容、数据、图表信息。用中文回复。"
            )
            vision_desc = call_model_vision("vision", prompt, non_text_urls)
            if vision_desc:
                idx_label = ",".join(str(i+1) for i in non_text_indices)
                result_parts.append(f"### [图片{idx_label} 图表描述]\n{vision_desc}")
                used_qwen = True
        except Exception as e:
            logger.warning(f"视觉模型降级失败: {e}")

    text = "\n\n".join(result_parts)
    # 如果没走 Qwen（全是 OCR），需要 DeepSeek 理解整理
    needs_understanding = not used_qwen
    return text, needs_understanding


def _extract_mixed(row: dict) -> str:
    """mixed: 分离文本和图片URL，文本保留 + 图片用 PaddleOCR（或视觉模型）（兼容旧调用）"""
    text, _ = _extract_mixed_with_meta(row)
    return text


def _extract_mixed_with_meta(row: dict) -> tuple:
    """mixed: 分离文本和图片URL，返回 (text, needs_understanding)"""
    text_content = row.get("text_content") or ""
    urls = _parse_image_urls(text_content)

    pure_text = re.sub(r"\[图片\d+\]\s*https?://\S+", "", text_content).strip()

    if not urls:
        return pure_text, True  # 纯文本，需要理解

    ocr_parts = []
    non_text_urls = []
    non_text_indices = []

    for idx, url in enumerate(urls):
        try:
            img_data = _download_file(url, timeout=30)
            ocr_text = _ocr_image_bytes(img_data)
            if len(ocr_text.strip()) >= 20:
                ocr_parts.append(f"### [图片{idx+1} OCR文字]\n{ocr_text}")
            else:
                non_text_urls.append(url)
                non_text_indices.append(idx)
        except Exception as e:
            logger.warning(f"PaddleOCR 图片失败: {e}")
            non_text_urls.append(url)
            non_text_indices.append(idx)

    img_parts = ocr_parts[:]
    used_qwen = False

    if non_text_urls:
        try:
            from utils.model_router import call_model_vision
            prompt = (
                f"这是关于'{row.get('title', '股票分析')}'的配图。"
                "请详细描述图片中的所有内容、数据、图表信息。用中文回复。"
            )
            vision_desc = call_model_vision("vision", prompt, non_text_urls)
            if vision_desc:
                idx_label = ",".join(str(i+1) for i in non_text_indices)
                img_parts.append(f"### [图片{idx_label} 图表描述]\n{vision_desc}")
                used_qwen = True
        except Exception as e:
            logger.warning(f"视觉模型降级失败: {e}")

    if not img_parts:
        return pure_text, True  # 图片全失败，纯文本

    text = f"{pure_text}\n\n--- 图片内容 ---\n" + "\n\n".join(img_parts)
    needs_understanding = not used_qwen
    return text, needs_understanding


def _get_zsxq_auth() -> dict:
    """获取知识星球认证信息（cookie + 请求头）"""
    try:
        from utils.db_utils import get_config as _gc
        token = _gc("zsxq_cookie") or os.getenv("ZSXQ_COOKIE", "")
    except Exception:
        token = os.getenv("ZSXQ_COOKIE", "")
    if not token:
        try:
            from config import ZSXQ_COOKIE
            token = ZSXQ_COOKIE
        except Exception:
            pass
    return {
        "cookies": {"zsxq_access_token": token} if token else {},
        "headers": {
            "origin": "https://wx.zsxq.com",
            "referer": "https://wx.zsxq.com/",
            "user-agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/144.0.0.0 Safari/537.36"
            ),
        } if token else {},
    }


def _download_file(url: str, timeout: int = 60, cookies: dict = None, headers: dict = None) -> bytes:
    """下载文件，SSL 失败时降级重试。知识星球域名自动带认证。"""
    import requests
    from urllib.parse import urlparse
    _cookies = dict(cookies or {})
    _headers = dict(headers or {})
    # 知识星球文件自动注入认证
    host = urlparse(url).hostname or ""
    if ("zsxq.com" in host or "zqimg.com" in host) and not _cookies:
        auth = _get_zsxq_auth()
        _cookies.update(auth["cookies"])
        _headers.update(auth["headers"])
    try:
        resp = requests.get(url, timeout=timeout, cookies=_cookies or None, headers=_headers or None)
        resp.raise_for_status()
        return resp.content
    except requests.exceptions.SSLError:
        logger.warning(f"SSL 错误，降级重试: {url[:80]}")
        resp = requests.get(url, timeout=timeout, verify=False, cookies=_cookies or None, headers=_headers or None)
        resp.raise_for_status()
        return resp.content


def _extract_page_main_column(page) -> str:
    """从 pdfplumber page 中提取主栏文字，过滤侧边栏噪音。

    三层过滤：列裁剪 → 整行匹配删除 → 行内碎片清理
    """
    words = page.extract_words(keep_blank_chars=False, x_tolerance=3, y_tolerance=3)
    if not words:
        return _filter_sidebar_lines(page.extract_text() or "")

    page_width = page.width
    if not page_width or page_width < 100:
        return _filter_sidebar_lines(page.extract_text() or "")

    # 列裁剪：研报主栏占左侧 ~72%
    main_cutoff = page_width * 0.72
    right_words = [w for w in words if w["x0"] > main_cutoff]

    # 右侧文字量 < 5% 说明无侧边栏
    if len(right_words) < len(words) * 0.05:
        return _filter_sidebar_lines(page.extract_text() or "")

    # 有侧边栏，裁剪主栏
    try:
        cropped = page.crop((0, 0, main_cutoff + 10, page.height))
        text = cropped.extract_text()
        if text and len(text.strip()) > 50:
            return _filter_sidebar_lines(text)
    except Exception:
        pass

    return _filter_sidebar_lines(page.extract_text() or "")


# 整行匹配的侧边栏模式
_SIDEBAR_LINE_PATTERNS = [
    re.compile(r"^执业证书编号[：:]?\s*S\d+"),
    re.compile(r"^[A-Za-z0-9_.+-]+@[A-Za-z0-9-]+\.[A-Za-z]{2,}$"),
    re.compile(r"^0\d{2,3}[-‐]\d{7,8}$"),
    re.compile(r"^(绝对|相对)表现%"),
    re.compile(r"^沪深300%"),
    re.compile(r"^1周\s+1月\s+3月\s+12月"),
    re.compile(r"^报告发布日期"),
    re.compile(r"^国家/地区\s+中国"),
    re.compile(r"^行业\s+\S{2,6}$"),
    re.compile(r"^(买入|增持|中性|减持|卖出)[（(](首次|维持)[)）]$"),
    re.compile(r"有关分析师的申明"),
    re.compile(r"请阅读本证券研究报告"),
    re.compile(r"本报告禁止在中国内地发布"),
    re.compile(r"^\d{1,3}$"),  # 纯页码
]

# 行内碎片清理模式
_SIDEBAR_INLINE_PATTERNS = [
    (re.compile(r"\s+报告发布日期\s*\d{4}年\d{1,2}月\d{1,2}日"), ""),
    (re.compile(r"\s+执业证书编号[：:]?\s*S\d+\S*"), ""),
    (re.compile(r"\s+\S+@\S+\.\S{2,4}$"), ""),
    (re.compile(r"\s+0\d{2,3}[-‐]\d{7,8}$"), ""),
    (re.compile(r"\s+\d{1,2}周\s+\d{1,2}月\s+\d{1,2}月\s+\d{1,2}月$"), ""),
]


def _filter_sidebar_lines(text: str) -> str:
    """正则过滤侧边栏噪音：整行删除 + 行内碎片清理"""
    if not text:
        return text
    lines = text.split("\n")
    filtered = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            filtered.append("")
            continue
        # 整行匹配删除
        if any(p.search(stripped) for p in _SIDEBAR_LINE_PATTERNS):
            continue
        # 行内碎片清理
        cleaned = line
        for pattern, repl in _SIDEBAR_INLINE_PATTERNS:
            cleaned = pattern.sub(repl, cleaned)
        filtered.append(cleaned)
    return "\n".join(filtered)


def _extract_office(data: bytes, url: str = "") -> str:
    """提取 Office Open XML 格式（docx/pptx），通过 ZIP 文件头识别"""
    import io, zipfile

    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
        names = zf.namelist()
    except Exception as e:
        raise ValueError(f"无法解析 ZIP/Office 文件: {e}")

    # docx: word/document.xml
    if "word/document.xml" in names:
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            # 降级：直接解析 XML
            import xml.etree.ElementTree as ET
            xml_data = zf.read("word/document.xml")
            root = ET.fromstring(xml_data)
            ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
            texts = [node.text for node in root.iter(f"{ns}t") if node.text]
            return "\n".join(texts)

    # pptx: ppt/slides/slide*.xml
    slide_files = sorted(n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
    if slide_files:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(p for p in parts if p.strip())
        except ImportError:
            import xml.etree.ElementTree as ET
            parts = []
            for sf in slide_files:
                xml_data = zf.read(sf)
                root = ET.fromstring(xml_data)
                texts = [node.text for node in root.iter() if node.text and node.text.strip()]
                parts.extend(texts)
            return "\n".join(parts)

    raise ValueError(f"未知 Office 格式，ZIP 内容: {names[:5]}")


def _extract_pdf_page_structured(page) -> str:
    """从单页提取结构化富文本：表格→Markdown，标题检测，段落分隔。

    保留 _extract_page_main_column 的侧边栏过滤逻辑，
    在其基础上增加表格提取和标题层级标记。
    """
    import statistics

    page_parts = []

    # 1. 提取表格区域（bbox 列表，用于后续排除）
    tables = page.find_tables() or []
    table_bboxes = [t.bbox for t in tables]

    # 把表格转 Markdown（按 y 坐标排序）
    table_md_by_top = {}
    for table in tables:
        try:
            data = table.extract()
            if not data or len(data) < 2:
                continue
            md_lines = []
            header = data[0]
            md_lines.append("| " + " | ".join(str(c or "") for c in header) + " |")
            md_lines.append("|" + "|".join("---" for _ in header) + "|")
            for row in data[1:]:
                md_lines.append("| " + " | ".join(str(c or "") for c in row) + " |")
            table_md_by_top[table.bbox[1]] = "\n".join(md_lines)
        except Exception:
            pass

    # 2. 提取文字（排除表格区域）
    words = page.extract_words(keep_blank_chars=False, x_tolerance=3, y_tolerance=3)
    if not words and not table_md_by_top:
        return _filter_sidebar_lines(page.extract_text() or "")

    # 侧边栏裁剪（复用原逻辑）
    page_width = page.width
    main_cutoff = page_width * 0.72 if page_width and page_width >= 100 else None
    if main_cutoff:
        right_words = [w for w in words if w["x0"] > main_cutoff]
        if len(right_words) >= len(words) * 0.05:
            # 有侧边栏，只保留主栏文字
            words = [w for w in words if w["x0"] <= main_cutoff + 10]

    # 过滤掉落在表格区域内的文字
    def _in_table(w):
        wx, wy = (w["x0"] + w.get("x1", w["x0"])) / 2, (w["top"] + w.get("bottom", w["top"])) / 2
        for bbox in table_bboxes:
            if bbox[0] - 2 <= wx <= bbox[2] + 2 and bbox[1] - 2 <= wy <= bbox[3] + 2:
                return True
        return False

    text_words = [w for w in words if not _in_table(w)]

    # 3. 标题检测：按字号分布，大于 median * 1.3 且单行 → 标题
    sizes = [w.get("size") or w.get("height", 10) for w in text_words]
    median_size = statistics.median(sizes) if sizes else 10

    # 按 y 坐标聚行
    lines_by_y = {}
    for w in text_words:
        y_key = round(w["top"], 1)
        # 合并接近的 y 值
        merged = False
        for yk in list(lines_by_y.keys()):
            if abs(yk - y_key) < 3:
                lines_by_y[yk].append(w)
                merged = True
                break
        if not merged:
            lines_by_y[y_key] = [w]

    # 4. 组装输出：按 y 坐标排序，穿插表格
    all_y_items = []  # (y, type, content)

    for y_top, line_words in sorted(lines_by_y.items()):
        line_words.sort(key=lambda w: w["x0"])
        line_text = " ".join(w["text"] for w in line_words).strip()
        if not line_text:
            continue

        # 检测标题
        avg_size = sum(w.get("size") or w.get("height", 10) for w in line_words) / len(line_words)
        is_heading = (avg_size > median_size * 1.3 and len(line_text) < 80)

        if is_heading:
            # 大字号 → ##，更大 → #
            prefix = "##" if avg_size < median_size * 1.8 else "#"
            all_y_items.append((y_top, "heading", f"{prefix} {line_text}"))
        else:
            all_y_items.append((y_top, "text", line_text))

    for y_top, md in table_md_by_top.items():
        all_y_items.append((y_top, "table", md))

    all_y_items.sort(key=lambda x: x[0])

    # 组装
    result_lines = []
    prev_type = None
    for _, item_type, content in all_y_items:
        if item_type == "heading":
            result_lines.append(f"\n{content}")
        elif item_type == "table":
            result_lines.append(f"\n{content}\n")
        else:
            # 连续文本行合并为段落
            if prev_type == "text" and result_lines:
                result_lines[-1] += " " + content
            else:
                result_lines.append(content)
        prev_type = item_type

    raw_text = "\n\n".join(result_lines)
    return _filter_sidebar_lines(raw_text)


def _describe_pdf_page_chart(page_img_path: str, page_num: int, title: str) -> str:
    """用视觉模型描述 PDF 某页中的图表，无图表则返回空字符串"""
    try:
        from utils.model_router import call_model_vision
        prompt = (
            f"这是一份研报「{title}」的第{page_num}页截图。"
            "请详细描述页面中所有图表的内容，包括图表类型、数据、趋势、关键数字。"
            '如果没有图表只有文字，回复"无图表"。用中文回复。'
        )
        result = call_model_vision("vision", prompt, [page_img_path],
                                   max_tokens=2000, timeout=60)
        if not result or result.strip().startswith("无图表"):
            return ""
        return result.strip()
    except Exception as e:
        logger.warning(f"视觉模型描述第{page_num}页失败: {e}")
        return ""


def _extract_pdf(row: dict) -> str:
    """pdf: 下载 oss_url 并提取全文（兼容旧调用）"""
    text, _ = _extract_pdf_with_meta(row)
    return text


def _extract_pdf_with_meta(row: dict) -> tuple:
    """pdf: 下载 oss_url 并提取全文，返回 (text, needs_understanding)

    策略：
    - 检测文件头，PK 开头说明是 Office 格式（docx/pptx），转用 python-docx/python-pptx
    - 否则用 pdfplumber 结构化提取（标题+表格+段落 Markdown 标记）
    - 每页截图送视觉模型（Qwen-VL-Max）描述图表内容
    - 若文字密度过低（扫描件），改用 PaddleOCR → needs_understanding=True
    - 正常 PDF 有 Qwen 图表描述 → needs_understanding=False
    """
    oss_url = row.get("oss_url")
    if not oss_url:
        return "", False

    import pdfplumber

    data = _download_file(oss_url)

    # 检测文件头：PK 是 ZIP/Office Open XML
    if data[:2] == b'PK':
        return _extract_office(data, oss_url), True  # Office 文档，需要理解

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        # 1. 截图所有页面（用于视觉模型）
        page_images = {}  # page_num(1-based) -> tmp_path
        try:
            from pdf2image import convert_from_path
            imgs = convert_from_path(tmp_path, dpi=150)
            for i, img in enumerate(imgs):
                img_path = f"{tmp_path}_page{i+1}.png"
                img.save(img_path, "PNG")
                page_images[i + 1] = img_path

            # 保存前2页低分辨率缩略图（供预览用，链接过期后仍可查看）
            doc_id = row.get("id")
            if doc_id:
                import pathlib
                thumb_dir = pathlib.Path(__file__).parent.parent / "static" / "uploads"
                thumb_dir.mkdir(parents=True, exist_ok=True)
                for pi in range(min(2, len(imgs))):
                    thumb = imgs[pi].copy()
                    thumb.thumbnail((800, 1200))
                    thumb_path = thumb_dir / f"thumb_{doc_id}_p{pi+1}.jpg"
                    thumb.save(str(thumb_path), "JPEG", quality=70)
        except Exception as e:
            logger.warning(f"PDF 截图失败，跳过图表提取: {e}")

        # 2. pdfplumber 提取文字+表格 + 视觉模型描述图表
        texts = []
        page_count = 0
        title = row.get("title") or "研报"

        with pdfplumber.open(tmp_path) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                page_num = i + 1
                # 文字+表格提取
                try:
                    page_text = _extract_pdf_page_structured(page)
                except Exception:
                    page_text = _extract_page_main_column(page)

                # 视觉模型描述图表
                chart_desc = ""
                if page_num in page_images:
                    chart_desc = _describe_pdf_page_chart(
                        page_images[page_num], page_num, title)

                # 拼合
                if page_text and chart_desc:
                    texts.append(f"{page_text}\n\n### [第{page_num}页 图表描述]\n{chart_desc}")
                elif page_text:
                    texts.append(page_text)
                elif chart_desc:
                    texts.append(f"### [第{page_num}页 图表描述]\n{chart_desc}")

        full_text = "\n\n".join(texts)
        char_density = len(full_text) / max(page_count, 1)

        # 每页平均字符数 < 100，判定为扫描件，改用 OCR
        if char_density < 100:
            logger.info(f"PDF 文字密度低（{char_density:.0f} 字/页），切换 PaddleOCR")
            try:
                ocr_text = _ocr_pdf_with_paddle(tmp_path)
                return ocr_text, True  # 扫描件 OCR，没过 Qwen，需要理解
            except Exception as e:
                logger.warning(f"PaddleOCR 失败，返回 pdfplumber 结果: {e}")

        # 正常 PDF：pdfplumber + Qwen 图表描述，轻度清洗即可
        return full_text, False
    finally:
        # 清理临时文件
        os.unlink(tmp_path)
        for img_path in page_images.values():
            try:
                os.unlink(img_path)
            except OSError:
                pass



def _extract_audio(row: dict) -> str:
    """audio: 下载 oss_url 并用 Groq/Gemini 转写"""
    oss_url = row.get("oss_url")
    if not oss_url:
        return ""

    import time

    data = _download_file(oss_url, timeout=120)

    # 推断文件扩展名
    ext = ".mp3"
    if ".wav" in oss_url:
        ext = ".wav"
    elif ".m4a" in oss_url:
        ext = ".m4a"

    mime_map = {".mp3": "audio/mpeg", ".wav": "audio/wav", ".m4a": "audio/mp4"}
    mime_type = mime_map.get(ext, "audio/mpeg")

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        cfg = _get_audio_config()

        # FunASR 本地优先
        if cfg["provider"] == "funasr":
            try:
                return _transcribe_funasr(tmp_path, cfg["model_name"] or "iic/SenseVoiceSmall")
            except Exception as e:
                logger.warning(f"FunASR 失败，降级到 Groq: {e}")

        groq_key = cfg["groq_key"]
        gemini_key = cfg["gemini_key"]

        if groq_key:
            try:
                return _transcribe_groq(tmp_path, groq_key)
            except Exception as e:
                if "rate_limit" in str(e) or "429" in str(e):
                    logger.warning("Groq 速率限制，切换 Gemini...")
                else:
                    raise

        if gemini_key:
            try:
                return _transcribe_gemini(data, gemini_key)
            except Exception as e:
                if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                    logger.warning("Gemini 也受限...")
                else:
                    raise

        if groq_key:
            logger.warning("两个 API 都受限，等待 10 分钟后重试 Groq...")
            time.sleep(620)
            return _transcribe_groq(tmp_path, groq_key)

        raise RuntimeError("未配置任何音频转写方案（funasr/groq/gemini）")
    finally:
        os.unlink(tmp_path)


def _parse_image_urls(text: str) -> list:
    """从 text_content 中解析 [图片N] url 格式的图片URL"""
    pattern = re.compile(r"\[图片\d+\]\s*(https?://\S+)")
    return pattern.findall(text)


# ==================== PaddleOCR ====================

_paddle_ocr_instance = None

def _get_paddle_ocr():
    """懒加载 PaddleOCR 实例（单例，避免重复初始化）"""
    global _paddle_ocr_instance
    if _paddle_ocr_instance is None:
        from paddleocr import PaddleOCR
        _paddle_ocr_instance = PaddleOCR(lang="ch")
    return _paddle_ocr_instance


def _ocr_image_bytes(image_bytes: bytes) -> str:
    """对图片字节数据做 OCR，返回识别文字"""
    import numpy as np
    from PIL import Image
    import io
    ocr = _get_paddle_ocr()
    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    img_array = np.array(img)
    result = ocr.predict(img_array)
    if not result:
        return ""
    # PaddleOCR 3.x: result[0]['rec_texts'] 是识别文字列表
    texts = result[0].get("rec_texts", []) if isinstance(result[0], dict) else []
    return "\n".join(t for t in texts if t)


def _ocr_pdf_with_paddle(pdf_path: str) -> str:
    """用 PaddleOCR 逐页识别 PDF（适用于扫描件）"""
    from pdf2image import convert_from_path
    import io
    texts = []
    pages = convert_from_path(pdf_path, dpi=200)
    for page_img in pages:
        buf = io.BytesIO()
        page_img.save(buf, format="PNG")
        page_text = _ocr_image_bytes(buf.getvalue())
        if page_text:
            texts.append(page_text)
    return "\n\n".join(texts)




def push_to_raw_items(limit: int = 50) -> dict:
    """将已提取的 source_documents 灌入 raw_items

    Args:
        limit: 每次处理数
    Returns:
        {"total": N, "pushed": M, "skipped": K}
    """
    rows = execute_cloud_query(
        """SELECT id, doc_type, file_type, title, extracted_text,
                  publish_date, source
           FROM source_documents
           WHERE extract_status IN ('extracted','ready_to_pipe','done') AND raw_item_id IS NULL
           LIMIT %s""",
        [limit],
    )

    # 获取 source_doc 的 source_id
    src_rows = execute_cloud_query(
        "SELECT id FROM data_sources WHERE name='source_doc'"
    )
    if not src_rows:
        raise RuntimeError("data_sources 中未找到 source_doc，请先运行 init_db")
    source_id = src_rows[0]["id"]

    total = len(rows)
    pushed = 0
    skipped = 0

    for row in rows:
        external_id = f"sd_{row['id']}"
        # 去重检查
        existing = execute_cloud_query(
            "SELECT id FROM raw_items WHERE source_id=%s AND external_id=%s",
            [source_id, external_id],
        )
        if existing:
            skipped += 1
            continue

        item_type = "report" if row["doc_type"] == "report" else "news"
        meta = json.dumps({
            "file_type": row["file_type"],
            "source_doc_id": row["id"],
            "original_source": row.get("source") or "",
            "sub_source": row.get("source") or "source_doc",
        }, ensure_ascii=False)

        raw_id = execute_cloud_insert(
            """INSERT INTO raw_items
               (source_id, external_id, title, content, published_at, item_type, meta_json)
               VALUES (%s, %s, %s, %s, %s, %s, %s)""",
            [source_id, external_id, row["title"],
             row["extracted_text"], row.get("publish_date"),
             item_type, meta],
        )
        # 回写 raw_item_id
        execute_cloud_insert(
            "UPDATE source_documents SET raw_item_id=%s WHERE id=%s",
            [raw_id, row["id"]],
        )
        pushed += 1

    logger.info(f"灌入raw_items完成: total={total}, pushed={pushed}, skipped={skipped}")
    return {"total": total, "pushed": pushed, "skipped": skipped}


# ==================== 灌入 extracted_texts（新管线）====================

def push_to_extracted_texts(limit: int = 50, on_progress=None) -> dict:
    """将已提取的 source_documents 灌入 extracted_texts（新管线）

    Args:
        limit: 每次处理数
        on_progress(done, total, row_id): 每完成一条时回调
    Returns:
        {"total": N, "pushed": M, "skipped": K, "failed": F}
    """
    from ingestion.base_source import _check_text_quality

    rows = execute_cloud_query(
        """SELECT id, doc_type, file_type, title, extracted_text,
                  publish_date, source
           FROM source_documents
           WHERE extract_status IN ('extracted','ready_to_pipe','done','remix')
             AND id NOT IN (
               SELECT source_doc_id FROM extracted_texts
               WHERE source_doc_id IS NOT NULL
             )
           LIMIT %s""",
        [limit],
    )

    total = len(rows)
    pushed = 0
    skipped = 0
    failed = 0

    if on_progress:
        on_progress(0, total, None)

    if not rows:
        return {"total": 0, "pushed": 0, "skipped": 0, "failed": 0}

    # 批量查已存在的 source_ref，避免逐条 SELECT
    source_refs = [f"sd_{row['id']}" for row in rows]
    placeholders = ",".join(["%s"] * len(source_refs))
    existing_refs = set()
    try:
        existing = execute_cloud_query(
            f"SELECT source_ref FROM extracted_texts WHERE source_ref IN ({placeholders})",
            source_refs,
        )
        existing_refs = {r["source_ref"] for r in (existing or [])}
    except Exception:
        pass

    for i, row in enumerate(rows):
        try:
            full_text = (row.get("extracted_text") or "").strip()
            if not full_text:
                skipped += 1
                if on_progress:
                    on_progress(i + 1, total, row["id"])
                continue

            source_ref = f"sd_{row['id']}"
            if source_ref in existing_refs:
                skipped += 1
                if on_progress:
                    on_progress(i + 1, total, row["id"])
                continue

            quality = _check_text_quality(full_text)
            source_name = row.get("source") or "source_doc"
            source_format = _file_type_to_format(row.get("file_type", "txt"))

            execute_cloud_insert(
                """INSERT INTO extracted_texts
                   (source, source_format, publish_time, full_text,
                    source_doc_id, source_ref, extract_quality, semantic_clean_status)
                   VALUES (%s, %s, %s, %s, %s, %s, %s, 'done')""",
                [source_name, source_format, row.get("publish_date"),
                 full_text, row["id"], source_ref, quality],
            )
            pushed += 1
        except Exception as e:
            logger.error(f"灌入 extracted_texts 失败 id={row['id']}: {e}")
            failed += 1
        if on_progress:
            on_progress(i + 1, total, row["id"])

    logger.info(f"灌入extracted_texts完成: total={total}, pushed={pushed}, "
                f"skipped={skipped}, failed={failed}")
    return {"total": total, "pushed": pushed, "skipped": skipped, "failed": failed}


def _file_type_to_format(file_type: str) -> str:
    """将 source_documents.file_type 映射为 extracted_texts.source_format"""
    mapping = {
        "txt": "text",
        "pdf": "pdf",
        "image": "image",
        "mixed": "mixed",
        "mp3": "audio",
        "audio": "audio",
        "xlsx": "xlsx",
        "xls": "xlsx",
    }
    return mapping.get(file_type, "text")


def extract_by_ids(doc_ids: list[int]) -> dict:
    """按指定 ID 列表提取源文档（定向提取）

    Args:
        doc_ids: source_documents.id 列表
    Returns:
        {"total": N, "success": M, "failed": K, "skipped": S}
    """
    if not doc_ids:
        return {"total": 0, "success": 0, "failed": 0, "skipped": 0}

    placeholders = ",".join(["%s"] * len(doc_ids))
    rows = execute_cloud_query(
        f"""SELECT id, doc_type, file_type, title, text_content, oss_url, extract_status
            FROM source_documents WHERE id IN ({placeholders})""",
        doc_ids,
    )
    total = len(rows)
    success = 0
    failed = 0
    skipped = 0

    for i, row in enumerate(rows):
        try:
            extracted = _extract_and_clean_single(row)
            from config.doc_types import classify_doc_type
            new_doc_type = classify_doc_type(
                row.get("title") or "",
                (extracted or "")[:200],
            )
            execute_cloud_insert(
                """UPDATE source_documents
                   SET extracted_text=%s, extract_status='extracted', doc_type=%s
                   WHERE id=%s""",
                [extracted, new_doc_type, row["id"]],
            )
            success += 1
        except Exception as e:
            err_str = str(e)
            if "401" in err_str or "403" in err_str or "Unauthorized" in err_str or "Forbidden" in err_str:
                logger.warning(f"跳过 id={row['id']} (URL 过期/无权限): {err_str[:100]}")
                execute_cloud_insert(
                    "UPDATE source_documents SET extract_status='skipped' WHERE id=%s",
                    [row["id"]],
                )
            else:
                logger.error(f"定向提取失败 id={row['id']}: {e}")
                execute_cloud_insert(
                    "UPDATE source_documents SET extract_status='failed' WHERE id=%s",
                    [row["id"]],
                )
            failed += 1

    logger.info(f"定向提取完成: total={total}, success={success}, failed={failed}, skipped={skipped}")
    return {"total": total, "success": success, "failed": failed, "skipped": skipped}


def push_to_extracted_texts_by_ids(doc_ids: list[int]) -> dict:
    """将指定 ID 的已提取文档灌入 extracted_texts（定向推入管线）

    Args:
        doc_ids: source_documents.id 列表（需 extract_status='done'）
    Returns:
        {"total": N, "pushed": M, "skipped": K, "failed": F}
    """
    from ingestion.base_source import _check_text_quality

    if not doc_ids:
        return {"total": 0, "pushed": 0, "skipped": 0, "failed": 0}

    placeholders = ",".join(["%s"] * len(doc_ids))
    rows = execute_cloud_query(
        f"""SELECT id, doc_type, file_type, title, extracted_text, publish_date, source
            FROM source_documents
            WHERE id IN ({placeholders}) AND extract_status IN ('extracted','ready_to_pipe','done','remix')""",
        doc_ids,
    )

    total = len(rows)
    pushed = 0
    skipped = 0
    failed = 0

    for row in rows:
        try:
            full_text = (row.get("extracted_text") or "").strip()
            if not full_text:
                skipped += 1
                continue

            source_ref = f"sd_{row['id']}"
            existing = execute_cloud_query(
                "SELECT id FROM extracted_texts WHERE source_ref=%s",
                [source_ref],
            )
            if existing:
                skipped += 1
                continue

            # 注意：extracted_text 已在 _extract_and_clean_single 阶段完成清洗
            # 此处直接入管线，不再重复清洗
            cleaned_text = full_text

            quality = _check_text_quality(cleaned_text)
            source_name = row.get("source") or "source_doc"
            source_format = _file_type_to_format(row.get("file_type", "txt"))

            execute_cloud_insert(
                """INSERT INTO extracted_texts
                   (source, source_format, publish_time, full_text,
                    source_doc_id, source_ref, extract_quality, semantic_clean_status)
                   VALUES (%s, %s, %s, %s, %s, %s, %s, 'done')""",
                [source_name, source_format, row.get("publish_date"),
                 cleaned_text, row["id"], source_ref, quality],
            )
            pushed += 1
        except Exception as e:
            logger.error(f"定向灌入 extracted_texts 失败 id={row['id']}: {e}")
            failed += 1

    logger.info(f"定向灌入extracted_texts完成: total={total}, pushed={pushed}, "
                f"skipped={skipped}, failed={failed}")
    return {"total": total, "pushed": pushed, "skipped": skipped, "failed": failed}
